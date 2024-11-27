import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import PyPDF2
import io

# Gemini API 키 설정
genai.configure(api_key='your_api_key')

# Gemini 모델 설정
model = genai.GenerativeModel('gemini-pro')

def create_presentation(title, content):
    try:
        # 프레젠테이션 생성
        prs = Presentation()
        
        # 제목 슬라이드
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # 내용을 9000자 단위로 분할
        def split_content(text, max_length=9000):
            chunks = []
            current_chunk = ""
            sentences = text.split('. ')
            
            for sentence in sentences:
                if len(current_chunk) + len(sentence) < max_length:
                    current_chunk += sentence + '. '
                else:
                    chunks.append(current_chunk)
                    current_chunk = sentence + '. '
            
            if current_chunk:
                chunks.append(current_chunk)
            return chunks
        
        content_chunks = split_content(content)
        
        # 각 청크에 대해 새로운 슬라이드 생성
        for chunk in content_chunks:
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            
            # 제목 없이 내용만 추가
            body_shape = slide.shapes.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.text = chunk
            
            # 텍스트 서식 설정
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(12)
        
        # 메모리에 저장
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        
        return pptx_io
    
    except Exception as e:
        st.error(f"프레젠테이션 생성 중 오류 발생: {str(e)}")
        return None

def extract_text_from_pdf(pdf_file):
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text
    except Exception as e:
        st.error(f"PDF 텍스트 추출 중 오류 발생: {str(e)}")
        return None

def main():
    st.title("PDF to PPT 변환기")
    
    uploaded_file = st.file_uploader("PDF 파일을 업로드하세요", type=['pdf'])
    
    if uploaded_file is not None:
        # PDF에서 텍스트 추출
        pdf_text = extract_text_from_pdf(uploaded_file)
        
        if pdf_text:
            try:
                # Gemini를 사용하여 내용 요약
                prompt = f"다음 텍스트를 PPT 형식으로 요약해주세요:\n\n{pdf_text}"
                response = model.generate_content(prompt)
                summary = response.text
                
                # PPT 생성
                title = "PDF 요약"
                pptx_io = create_presentation(title, summary)
                
                if pptx_io:
                    # 다운로드 버튼 생성
                    st.download_button(
                        label="PPT 다운로드",
                        data=pptx_io.getvalue(),
                        file_name="summary.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
            
            except Exception as e:
                st.error(f"처리 중 오류 발생: {str(e)}")

if __name__ == "__main__":
    main()
